"""File content extraction utilities for GDPVal grading.

Extracts text content from various file formats (PDF, XLSX, DOCX, PPTX, ZIP, plain text)
for use by the LLM-as-Judge grader.
"""

from __future__ import annotations

import logging
import tempfile
import zipfile
from pathlib import Path

logger = logging.getLogger(__name__)

MAX_CHARS_PER_FILE = 50_000


def extract_text(path: Path) -> str:
    """Extract text content from a file, dispatching by extension.

    Returns plain text representation of the file content,
    truncated to MAX_CHARS_PER_FILE characters.
    """
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            text = _extract_pdf(path)
        elif suffix in (".xlsx", ".xls", ".xlsm"):
            text = _extract_xlsx(path)
        elif suffix in (".docx", ".doc"):
            text = _extract_docx(path)
        elif suffix in (".pptx", ".ppt"):
            text = _extract_pptx(path)
        elif suffix == ".zip":
            text = _extract_zip(path)
        else:
            text = _extract_text_file(path)
    except Exception as exc:
        logger.warning("Failed to extract text from %s: %s", path, exc)
        return f"[Could not extract content from {path.name}: {exc}]"

    if len(text) > MAX_CHARS_PER_FILE:
        text = text[:MAX_CHARS_PER_FILE] + f"\n\n[... truncated at {MAX_CHARS_PER_FILE} characters]"
    return text


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"--- Page {i + 1} ---\n{text}")
    return "\n\n".join(pages)


def _extract_xlsx(path: Path) -> str:
    """Extract text from Excel using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append("\t".join(cells))
        sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets)


def _extract_docx(path: Path) -> str:
    """Extract text from Word document using python-docx."""
    from docx import Document

    doc = Document(path)
    parts = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    for table in doc.tables:
        table_rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            table_rows.append("\t".join(cells))
        parts.append("[Table]\n" + "\n".join(table_rows))

    return "\n\n".join(parts)


def _extract_pptx(path: Path) -> str:
    """Extract text from PowerPoint using python-pptx."""
    from pptx import Presentation

    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
        if texts:
            slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_zip(path: Path) -> str:
    """Extract text from ZIP by unpacking and recursively processing files."""
    parts = []
    with tempfile.TemporaryDirectory() as tmpdir:
        with zipfile.ZipFile(path, "r") as zf:
            zf.extractall(tmpdir)
        tmp_path = Path(tmpdir)
        for file in sorted(tmp_path.rglob("*")):
            if file.is_file() and not file.name.startswith("."):
                rel = file.relative_to(tmp_path)
                text = extract_text(file)
                parts.append(f"--- {rel} ---\n{text}")
    return "\n\n".join(parts)


def _extract_text_file(path: Path) -> str:
    """Read plain text files (txt, md, csv, json, py, etc.)."""
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        try:
            return path.read_text(encoding="latin-1")
        except Exception:
            return f"[Binary file: {path.name}]"
